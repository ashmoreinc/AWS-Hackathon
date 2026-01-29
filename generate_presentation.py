from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        top = Inches(4)
        txBox = slide.shapes.add_textbox(left, top, width, Inches(1))
        tf = txBox.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 0
        p.space_before = Pt(8)

# Slide 1: Title
add_title_slide("Offer Management System", "Dynamic Discount Optimization")

# Slide 2: Challenge Overview
add_content_slide("Hackathon Challenge", [
    "Create a dynamic system that optimizes discount display in real-time",
    "Based on: User profile, redemption history, inventory, merchant priorities",
    "Balance high-value offers with diverse deal mix",
    "Adapt to: Trending categories, expiring promotions, environmental factors"
])

# Slide 3: Core Data Models
add_content_slide("Core Data Models", [
    "Offer Data: JSON file with all available offers",
    "User Data: JSON file with user profiles and preferences",
    "Offer Types: online, in-store, local",
    "User Connection Types:",
    "  • wifi - User at home",
    "  • mobile - User on the go"
])

# Slide 4: Pre-filtering
add_content_slide("Pre-filtering Rules", [
    "Inventory Check:",
    "  • Only include offers where inventory_count > 0",
    "  • Offers with no inventory field are always included",
    "",
    "This ensures users only see available offers"
])

# Slide 5: Base Scoring (1/2)
add_content_slide("Base Contextual Scoring (0-4 points)", [
    "Connection Match: +1 point",
    "  • WiFi + online offer = +1",
    "  • Mobile + non-online offer = +1",
    "",
    "Merchant Priority Boost: +1 point",
    "  • If offer.boost == true"
])

# Slide 6: Base Scoring (2/2)
add_content_slide("Base Scoring Continued", [
    "Urgency (Expiry): 0-1 point",
    "  • < 7 days = +1 point (critical)",
    "  • 7-14 days = +0.5 point (urgent)",
    "  • > 14 days = 0 points",
    "",
    "Commission Value: 0-1 point",
    "  • 0-1 = 0 points (low)",
    "  • 2-3 = +0.5 points (medium)",
    "  • 4-5 = +1 point (high)"
])

# Slide 7: Advanced Scoring
add_content_slide("Advanced Scoring Features", [
    "Inventory Scarcity: 0-0.3 points",
    "  • Low inventory (1-10) = +0.3 points",
    "  • High inventory (>100) = 0 points",
    "",
    "Personalization Boost: 0-1 point",
    "  • AWS Personalize score added directly",
    "  • Only for recommended offers",
    "",
    "Lunch Time Boost: +3 points",
    "  • 12:00 PM - 2:00 PM + offer.lunch == true"
])

# Slide 8: Penalties & Constraints
add_content_slide("Penalties & Diversity", [
    "Gift Card Penalty: -0.5 points",
    "  • Lower purchase intent for gift cards",
    "",
    "Coffee Time Penalty: 0 to -1 points (linear)",
    "  • Deprioritize coffee offers after 2:00 PM",
    "  • At 2PM = 0, at 6PM = -0.4, at midnight = -1.0",
    "  • Detected by title: coffee, café, espresso, latte, etc.",
    "",
    "Diversity Constraint:",
    "  • Maximum 60% (12 out of 20) of any single offerType"
])

# Slide 9: Scoring Flow
add_content_slide("Scoring Flow Process", [
    "1. Pre-filter: Remove zero inventory offers",
    "2. Calculate base score: Connection + Boost + Expiry + Commission",
    "3. Add inventory scarcity bonus",
    "4. Add personalization boost (AWS Personalize)",
    "5. Add lunch time boost if applicable",
    "6. Apply gift card penalty if applicable",
    "7. Apply coffee time penalty if after 2PM",
    "8. Sort by final score (descending)",
    "9. Apply diversity constraint: Select top 20"
])

# Slide 10: Example - WiFi User
add_content_slide("Example: WiFi User", [
    "User: USER001 (connection_type: wifi)",
    "",
    "Amazon Prime (online, boost, commission: 5)",
    "  → Score: +1 (wifi) +1 (boost) +1 (commission) = 3 ✓",
    "",
    "ASOS Sale (online, commission: 3)",
    "  → Score: +1 (wifi) +0.5 (commission) = 1.5 ✓",
    "",
    "Starbucks (in-store, commission: 2)",
    "  → Score: +0.5 (commission) = 0.5"
])

# Slide 11: Example - Mobile User
add_content_slide("Example: Mobile User", [
    "User: USER002 (connection_type: mobile)",
    "",
    "Starbucks (in-store, boost, commission: 4)",
    "  → Score: +1 (mobile) +1 (boost) +1 (commission) = 3 ✓",
    "",
    "Tesco (local, commission: 3)",
    "  → Score: +1 (mobile) +0.5 (commission) = 1.5 ✓",
    "",
    "Amazon Prime (online, commission: 5)",
    "  → Score: +1 (commission) = 1"
])

# Slide 12: API Behavior
add_content_slide("API Features", [
    "Connection Type Override:",
    "  • Override user's stored connection_type",
    "  • Override current time for testing",
    "",
    "Response Format:",
    "  • Top 20 offers sorted by final_score",
    "  • User's connection_type",
    "  • base_score: Business rules score (0-4.3)",
    "  • personalize_score: AWS Personalize (0-1)",
    "  • final_score: base_score + personalize_score"
])

# Slide 13: System Constraints
add_content_slide("System Limits & Constraints", [
    "Maximum 20 offers per request",
    "  • Fast response time",
    "  • Prevents decision paralysis",
    "",
    "Lambda Timeout: 30 seconds",
    "Kinesis Retention: 24 hours",
    "DynamoDB: Pay-per-request (auto-scales)",
    "User activity stored indefinitely"
])

# Slide 14: Summary
add_content_slide("Summary", [
    "Intelligent ranking system balancing multiple factors",
    "Real-time personalization with AWS Personalize",
    "Context-aware (connection type, time, inventory)",
    "Diversity constraints ensure balanced offer mix",
    "Scalable serverless architecture",
    "Maximum 20 highly relevant offers per user"
])

prs.save('/Users/khaldounchahde/WebstormProjects/untitled/Offer_Management_System_Presentation.pptx')
print("✓ Presentation created: Offer_Management_System_Presentation.pptx")
